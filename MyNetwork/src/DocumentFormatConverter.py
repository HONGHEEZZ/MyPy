from PyQt5.QtWidgets import *
from PyQt5.QtCore import *
from PyQt5.QtGui import *
from PyQt5 import uic
import os
import win32com.client
import time

#for ppt 텍스트 변경 : pip install python-pptx
from pptx import Presentation






word = win32com.client.Dispatch("Word.Application")

my_dlg = uic.loadUiType("../ui/DocumentFormatConverter.ui")[0]





class CDocumentFormatConverter(QDialog, my_dlg):
    def __init__(self,parent=None):
        print("parent:", parent)
        super().__init__(parent)
        self.setupUi(self)
        #loadUi("test2",self)

        self.initUI()

        self.btnView.clicked.connect(self.btnView_Click)
        self.btnConvert.clicked.connect(self.btnConvert_Click)

        self.btnPPT.clicked.connect(self.btnPPT_Click)

        self.txtSourceDir.setText(r"D:\test")
        self.txtTargetDir.setText(r"D:\test\output")

    def initUI(self):
        self.setWindowTitle('DocumentFormatConverter')

        self.setMinimumSize(200, 200)

        self.tblFiles.setColumnCount(2)
        self.tblFiles.setHorizontalHeaderLabels(["FileName", "Directory"])

        self.tblFiles.setColumnWidth(0, 200)
        self.tblFiles.setColumnWidth(1, 400)

        self.progressBar.reset()

        self.cboOutput.addItem("0 : doc : Microsoft Office Word 97 - 2003 binary file format.")
        self.cboOutput.addItem("2 : txt : Microsoft Windows text format.")
        self.cboOutput.addItem("8 : html : Standard HTML format.")
        self.cboOutput.addItem("16: docx : Word default document file format. DOCX")
        self.cboOutput.addItem("17: pdf : PDF format.")

        #self.tblFiles.rowcount = 1

    def btnView_Click(self):
        myDir = self.txtSourceDir.text()
        myExtend  = self.txtExtend.text()
        if not myExtend.startswith('.'):
            myExtend = '.' + myExtend

        list = []
        search(list, myDir, myExtend)
        print(list)

        self.tblFiles.setRowCount(len(list))


        for i, file in enumerate(list):
            self.tblFiles.setItem(i, 0, QTableWidgetItem(file[1]))
            self.tblFiles.setItem(i, 1, QTableWidgetItem(file[0]))

    def btnPPT_Click(self):

        out_Dir = self.txtTargetDir.text()

        cnt = self.tblFiles.rowCount()

        # 상태 진행바
        self.progressBar.setRange(0, cnt)

        row = 0
        while True:
            cell = self.tblFiles.item(row, 1)
            dir_name = cell.text()

            cell = self.tblFiles.item(row, 0)
            file_name = cell.text()


            #def replaceTextPPT(dir_name, file_name, out_dir):
            replaceTextPPT(dir_name, file_name, out_Dir)
            row = row + 1
            # 상태 진행바
            self.progressBar.setValue(row)

    def btnConvert_Click(self):
        toFileFormat = self.cboOutput.currentText()
        lisfFileFormat = toFileFormat.split(':')
        toFileFormatNum = lisfFileFormat[0].strip()
        toFileFormatNum = int(toFileFormatNum)
        toFileFormatExt = lisfFileFormat[1].strip()
        print("********convertTo:",toFileFormat)

        out_Dir = self.txtTargetDir.text()

        cnt = self.tblFiles.rowCount()

        #상태 진행바
        self.progressBar.setRange(0,cnt)

        row = 0
        while True:

            
            cell = self.tblFiles.item(row, 1)
            dir_name = cell.text()

            cell = self.tblFiles.item(row, 0)
            file_name = cell.text()

            pdfToOtherFormat(dir_name, file_name, out_Dir, toFileFormatNum, toFileFormatExt)
            row = row + 1
            # 상태 진행바
            self.progressBar.setValue(row)
            

def replaceTextPPT(dir_name, file_name, out_dir):
    full_filename = os.path.join(dir_name, file_name)
    full_out_file_name = os.path.join(out_dir, file_name)


    print("----------------------------------------------------")
    print(full_filename)
    print("----------------------------------------------------")

    #file_name = r"d:\test\2.2.1.1 데이터 요건정의 및 입수프로세스구축-데이터관리.pptx"
    try:
        prs = Presentation(full_filename)
    except:
        print("**************에러발생")
        return


    result = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:

                # 문장이 .로 끝날때만 대상으로 함. --> 여기서 바꾸면 폰트가 바뀜
                if paragraph.text.strip().endswith("."):

                    # 이렇게 바꿔야 폰트가 안바뀜....
                    for run in paragraph.runs:
                        print(run.text)
                        if run.text.strip().endswith("."):
                            print(run.text)
                            run.text = run.text.strip()[:-1]
                        result.append(run.text)

    prs.save(full_out_file_name)

def search(list, dir_name, myExtend, findSubDir=True):

    file_names = os.listdir(dir_name)

    for file_name in file_names:
        full_filename = os.path.join(dir_name, file_name)
        #print(full_filename)
        if findSubDir and os.path.isdir(full_filename):
            search(list, full_filename, myExtend, findSubDir)

        else:
            ext = os.path.splitext(full_filename)[-1]
            if ext == myExtend:
                print('*************pdf 파일임...', full_filename)
                # pdfTohtml(dir_name, file_name)
                #wordTotxt(dir_name, file_name)
                list.append([dir_name, file_name])




def pdfToOtherFormat(dir_name, file_name, out_dir, toFileFormatNum, toFileFormatExt):
    word.Visible = True
    word.DisplayAlerts = False

    full_in_file_name = os.path.join(dir_name, file_name)
    print("* full_in_file_name : ", full_in_file_name)


    in_name = os.path.splitext(file_name)[0]
    in_ext = os.path.splitext(file_name)[-1]
    out_name = in_name


    full_out_file_name = os.path.join(out_dir, out_name) + '.' + toFileFormatExt

    # Open and copy HTML
    print('* PDF 파일 열기 시작----', full_in_file_name)
    doc = word.Documents.Add(full_in_file_name)
    word.Selection.WholeStory()
    word.Selection.Copy()
    doc.Close(SaveChanges=False)
    print('* PDF 파일 복사 후 닫기 종료----', full_in_file_name)

    # Open new document, paste HTML and save
    doc = word.Documents.Add()
    word.Selection.Paste()
    print('* PDF 파일 전체 선택 후 붙여넣기 종료----', full_in_file_name)



    print('* 다른이름으로 저장 시작----{}, {}'.format(full_out_file_name, toFileFormatNum))
    doc.SaveAs(full_out_file_name, FileFormat=toFileFormatNum)  # 8	Standard HTML format.


    from ctypes import windll
    if windll.user32.OpenClipboard(None):
        windll.user32.EmptyClipboard()
        windll.user32.CloseClipboard()
        print('* 클립보드 클리어', full_in_file_name)

    doc.Close(SaveChanges=False)
    print('* Word app 종료', full_in_file_name)


if __name__ == "__main__":

    import sys
    app = QApplication(sys.argv)
    ex = CDocumentFormatConverter()
    ex.show()
    sys.exit(app.exec_())



'''
Name	Value	Description
wdFormatDocument	0	Microsoft Office Word 97 - 2003 binary file format.
wdFormatDOSText	4	Microsoft DOS text format.
wdFormatDOSTextLineBreaks	5	Microsoft DOS text with line breaks preserved.
wdFormatEncodedText	7	Encoded text format.
wdFormatFilteredHTML	10	Filtered HTML format.
wdFormatFlatXML	19	Open XML file format saved as a single XML file.
wdFormatFlatXMLMacroEnabled	20	Open XML file format with macros enabled saved as a single XML file.
wdFormatFlatXMLTemplate	21	Open XML template format saved as a XML single file.
wdFormatFlatXMLTemplateMacroEnabled	22	Open XML template format with macros enabled saved as a single XML file.
wdFormatOpenDocumentText	23	OpenDocument Text format.
wdFormatHTML	8	Standard HTML format.
wdFormatRTF	6	Rich text format (RTF).
wdFormatStrictOpenXMLDocument	24	Strict Open XML document format.
wdFormatTemplate	1	Word template format.
wdFormatText	2	Microsoft Windows text format.
wdFormatTextLineBreaks	3	Windows text format with line breaks preserved.
wdFormatUnicodeText	7	Unicode text format.
wdFormatWebArchive	9	Web archive format.
wdFormatXML	11	Extensible Markup Language (XML) format.
wdFormatDocument97	0	Microsoft Word 97 document format.
wdFormatDocumentDefault	16	Word default document file format. For Word, this is the DOCX format.
wdFormatPDF	17	PDF format.
wdFormatTemplate97	1	Word 97 template format.
wdFormatXMLDocument	12	XML document format.
wdFormatXMLDocumentMacroEnabled	13	XML document format with macros enabled.
wdFormatXMLTemplate	14	XML template format.
wdFormatXMLTemplateMacroEnabled	15	XML template format with macros enabled.
wdFormatXPS	18	XPS format.
'''


'''


* 파이썬 ui 참고 자료 - wikidoc
- PyQt5 Tutorial - 파이썬으로 만드는 나만의 GUI 프로그램
- 초보자를 위한 Python GUI 프로그래밍 - PyQt5
- python_ Cheat Sheet

'''







'''
-------------------------------------------------------------------------------
* PpSaveAsFileType enumeration (PowerPoint)
-------------------------------------------------------------------------------


https://docs.microsoft.com/en-us/office/vba/api/powerpoint.ppsaveasfiletype




Constants that specify type of file to save as, passed to the SaveAs method of the Presentation object.

PPSAVEASFILETYPE ENUMERATION (POWERPOINT)
ppSaveAsAddIn	8
ppSaveAsAnimatedGIF	40
ppSaveAsBMP	19
ppSaveAsDefault	11
ppSaveAsEMF	23
ppSaveAsExternalConverter	64000
ppSaveAsGIF	16
ppSaveAsJPG	17
ppSaveAsMetaFile	15
ppSaveAsMP4	39
ppSaveAsOpenDocumentPresentation	35
ppSaveAsOpenXMLAddin	30
ppSaveAsOpenXMLPicturePresentation	36
ppSaveAsOpenXMLPresentation	24
ppSaveAsOpenXMLPresentationMacroEnabled	25
ppSaveAsOpenXMLShow	28
ppSaveAsOpenXMLShowMacroEnabled	29
ppSaveAsOpenXMLTemplate	26
ppSaveAsOpenXMLTemplateMacroEnabled	27
ppSaveAsOpenXMLTheme	31
ppSaveAsPDF	32
ppSaveAsPNG	18
ppSaveAsPresentation	1
ppSaveAsRTF	6
ppSaveAsShow	7
ppSaveAsStrictOpenXMLPresentation	38
ppSaveAsTemplate	5
ppSaveAsTIF	21
ppSaveAsWMV	37
ppSaveAsXMLPresentation	34
ppSaveAsXPS	33

'''